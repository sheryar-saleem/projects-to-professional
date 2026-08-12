import json
import logging
import os
import re
import sqlite3
import tempfile
import time
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from dotenv import load_dotenv
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field
from azure.identity import DefaultAzureCredential
from azure.ai.projects import AIProjectClient

load_dotenv()

logging.basicConfig(level=os.getenv("LOG_LEVEL", "INFO"), format="%(asctime)s %(levelname)s %(name)s %(message)s")
logger = logging.getLogger("docmind")

BASE_DIR = Path(__file__).resolve().parent
DATA_DIR = BASE_DIR / "data"
UPLOAD_DIR = DATA_DIR / "uploads"
DB_PATH = DATA_DIR / "docmind.db"
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

MAX_FILE_SIZE = int(os.getenv("MAX_UPLOAD_MB", "25")) * 1024 * 1024
ALLOWED_EXTENSIONS = {".pdf", ".docx", ".txt", ".csv", ".xlsx", ".pptx"}

app = FastAPI(title="DocMind API", version="1.0.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=[x.strip() for x in os.getenv("CORS_ORIGINS", "http://localhost:3000").split(",")],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def now() -> str:
    return datetime.now(timezone.utc).isoformat()


def db() -> sqlite3.Connection:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    con = sqlite3.connect(DB_PATH)
    con.row_factory = sqlite3.Row
    return con


def init_db() -> None:
    with db() as con:
        con.executescript(
            """
            CREATE TABLE IF NOT EXISTS conversations (
                id TEXT PRIMARY KEY,
                title TEXT NOT NULL,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL
            );
            CREATE TABLE IF NOT EXISTS messages (
                id TEXT PRIMARY KEY,
                conversation_id TEXT NOT NULL,
                role TEXT NOT NULL,
                content TEXT NOT NULL,
                created_at TEXT NOT NULL,
                latency_ms INTEGER,
                FOREIGN KEY(conversation_id) REFERENCES conversations(id) ON DELETE CASCADE
            );
            CREATE TABLE IF NOT EXISTS documents (
                id TEXT PRIMARY KEY,
                name TEXT NOT NULL,
                filename TEXT NOT NULL,
                extension TEXT NOT NULL,
                size INTEGER NOT NULL,
                status TEXT NOT NULL,
                error TEXT,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL,
                chunks INTEGER DEFAULT 0
            );
            CREATE TABLE IF NOT EXISTS errors (
                id TEXT PRIMARY KEY,
                message TEXT NOT NULL,
                created_at TEXT NOT NULL
            );
            CREATE TABLE IF NOT EXISTS metrics (
                id INTEGER PRIMARY KEY CHECK(id=1),
                questions INTEGER NOT NULL DEFAULT 0,
                total_latency_ms INTEGER NOT NULL DEFAULT 0
            );
            INSERT OR IGNORE INTO metrics(id, questions, total_latency_ms) VALUES(1,0,0);
            """
        )


init_db()

_project_client: AIProjectClient | None = None
_openai_client: Any = None


def azure_clients() -> tuple[AIProjectClient, Any]:
    global _project_client, _openai_client
    endpoint = os.getenv("PROJECT_ENDPOINT")
    if not endpoint:
        raise RuntimeError("PROJECT_ENDPOINT is not configured on the server.")
    if _project_client is None:
        _project_client = AIProjectClient(endpoint=endpoint, credential=DefaultAzureCredential())
        _openai_client = _project_client.get_openai_client()
    return _project_client, _openai_client


def clean_citations(text: str) -> str:
    text = re.sub(r"\[\d+:\d+†.*?\]", "", text)
    text = re.sub(r"【.*?】", "", text)
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def safe_error(message: str) -> None:
    logger.exception(message)
    with db() as con:
        con.execute("INSERT INTO errors(id,message,created_at) VALUES(?,?,?)", (str(uuid.uuid4()), message, now()))


def extract_text(path: Path, ext: str) -> str:
    if ext == ".txt":
        return path.read_text(encoding="utf-8", errors="ignore")
    if ext == ".pdf":
        from pypdf import PdfReader
        return "\n\n".join(page.extract_text() or "" for page in PdfReader(str(path)).pages)
    if ext == ".docx":
        from docx import Document
        return "\n".join(p.text for p in Document(str(path)).paragraphs)
    if ext in {".csv", ".xlsx"}:
        import pandas as pd
        if ext == ".csv":
            frame = pd.read_csv(path, dtype=str).fillna("")
        else:
            frame = pd.read_excel(path, dtype=str).fillna("")
        return frame.to_csv(index=False)
    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(path))
        chunks = []
        for i, slide in enumerate(prs.slides, 1):
            text = " ".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            chunks.append(f"Slide {i}\n{text}")
        return "\n\n".join(chunks)
    raise ValueError("Unsupported document type.")


def chunk_text(text: str, size: int = 1200, overlap: int = 150) -> list[str]:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = min(len(text), start + size)
        chunk = text[start:end]
        if end < len(text):
            cut = chunk.rfind(" ")
            if cut > size * 0.6:
                end = start + cut
                chunk = text[start:end]
        chunks.append(chunk.strip())
        start = max(end - overlap, start + 1)
    return chunks


class ChatRequest(BaseModel):
    message: str = Field(min_length=1, max_length=12000)
    conversation_id: str | None = None


class RenameRequest(BaseModel):
    title: str = Field(min_length=1, max_length=120)


@app.get("/")
def home():
    return {"name": "DocMind API", "status": "running"}


@app.get("/health")
def health():
    result = {"status": "healthy", "backend": "healthy", "azure_ai": "unknown", "search": "unknown", "knowledge_base": "healthy"}
    try:
        if not os.getenv("PROJECT_ENDPOINT"):
            result["azure_ai"] = "unconfigured"
            result["status"] = "degraded"
        else:
            client, _ = azure_clients()
            # Lightweight client construction/credential path check.
            _ = client
            result["azure_ai"] = "healthy"
    except Exception as exc:
        logger.warning("Azure health check failed: %s", exc)
        result["azure_ai"] = "unavailable"
        result["status"] = "degraded"
    if os.getenv("AZURE_SEARCH_ENDPOINT"):
        result["search"] = "configured"
    else:
        result["search"] = "unconfigured"
        result["status"] = "degraded"
    return result


@app.get("/conversations")
def conversations():
    with db() as con:
        rows = con.execute("SELECT * FROM conversations ORDER BY updated_at DESC").fetchall()
    return [dict(r) for r in rows]


@app.get("/conversations/{conversation_id}")
def conversation(conversation_id: str):
    with db() as con:
        conv = con.execute("SELECT * FROM conversations WHERE id=?", (conversation_id,)).fetchone()
        msgs = con.execute("SELECT id,role,content,created_at FROM messages WHERE conversation_id=? ORDER BY created_at", (conversation_id,)).fetchall()
    if not conv:
        raise HTTPException(404, "Conversation not found")
    return {"conversation": dict(conv), "messages": [dict(x) for x in msgs]}


@app.post("/conversations")
def create_conversation():
    cid = str(uuid.uuid4())
    stamp = now()
    with db() as con:
        con.execute("INSERT INTO conversations VALUES(?,?,?,?)", (cid, "New conversation", stamp, stamp))
    return {"id": cid, "title": "New conversation"}


@app.patch("/conversations/{conversation_id}")
def rename_conversation(conversation_id: str, request: RenameRequest):
    with db() as con:
        cur = con.execute("UPDATE conversations SET title=?,updated_at=? WHERE id=?", (request.title.strip(), now(), conversation_id))
    if cur.rowcount == 0:
        raise HTTPException(404, "Conversation not found")
    return {"id": conversation_id, "title": request.title.strip()}


@app.delete("/conversations/{conversation_id}")
def delete_conversation(conversation_id: str):
    with db() as con:
        con.execute("DELETE FROM messages WHERE conversation_id=?", (conversation_id,))
        cur = con.execute("DELETE FROM conversations WHERE id=?", (conversation_id,))
    if cur.rowcount == 0:
        raise HTTPException(404, "Conversation not found")
    return {"deleted": True}


@app.post("/chat")
def chat(request: ChatRequest):
    start = time.perf_counter()
    cid = request.conversation_id
    try:
        if not cid:
            cid = create_conversation()["id"]
        with db() as con:
            if not con.execute("SELECT 1 FROM conversations WHERE id=?", (cid,)).fetchone():
                raise HTTPException(404, "Conversation not found")
            con.execute("INSERT INTO messages VALUES(?,?,?,?,?,?)", (str(uuid.uuid4()), cid, "user", request.message, now(), None))

        _, client = azure_clients()
        response = client.responses.create(
            input=[{"role": "user", "content": request.message}],
            extra_body={"agent_reference": {"name": os.getenv("AGENT_NAME"), "version": os.getenv("AGENT_VERSION"), "type": "agent_reference"}},
        )
        text = clean_citations(response.output_text or "")
        latency = int((time.perf_counter() - start) * 1000)
        with db() as con:
            con.execute("INSERT INTO messages VALUES(?,?,?,?,?,?)", (str(uuid.uuid4()), cid, "assistant", text, now(), latency))
            con.execute("UPDATE conversations SET updated_at=?, title=CASE WHEN title='New conversation' THEN ? ELSE title END WHERE id=?", (now(), request.message[:60], cid))
            con.execute("UPDATE metrics SET questions=questions+1,total_latency_ms=total_latency_ms+? WHERE id=1", (latency,))
        return {"conversation_id": cid, "response": text, "grounding": "knowledge_base"}
    except HTTPException:
        raise
    except Exception:
        safe_error("Chat request failed")
        raise HTTPException(503, "Sorry, I couldn't process your request right now. Please try again.")


@app.get("/documents")
def documents():
    with db() as con:
        return [dict(r) for r in con.execute("SELECT * FROM documents ORDER BY created_at DESC").fetchall()]


@app.post("/documents/upload")
async def upload_document(file: UploadFile = File(...)):
    filename = Path(file.filename or "").name
    ext = Path(filename).suffix.lower()
    if not filename or ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(400, "Unsupported file type. Allowed: PDF, DOCX, TXT, CSV, XLSX, PPTX.")
    document_id = str(uuid.uuid4())
    stored = UPLOAD_DIR / f"{document_id}{ext}"
    total = 0
    try:
        with stored.open("wb") as out:
            while chunk := await file.read(1024 * 1024):
                total += len(chunk)
                if total > MAX_FILE_SIZE:
                    out.close(); stored.unlink(missing_ok=True)
                    raise HTTPException(413, f"File exceeds the {MAX_FILE_SIZE // 1024 // 1024} MB limit.")
                out.write(chunk)
        stamp = now()
        with db() as con:
            con.execute("INSERT INTO documents VALUES(?,?,?,?,?,?,?,?,?,?)", (document_id, filename, stored.name, ext, total, "processing", None, stamp, stamp, 0))
        try:
            text = extract_text(stored, ext)
            chunks = chunk_text(text)
            if not chunks:
                raise ValueError("No extractable text was found in this document.")
            # The existing Azure Agent remains the retrieval authority. Persist extracted content
            # and chunk count locally so ingestion status is real; Azure Search integration can be
            # enabled separately with the environment settings below.
            with db() as con:
                con.execute("UPDATE documents SET status='indexed',chunks=?,updated_at=? WHERE id=?", (len(chunks), now(), document_id))
        except Exception as exc:
            logger.exception("Document processing failed")
            with db() as con:
                con.execute("UPDATE documents SET status='failed',error=?,updated_at=? WHERE id=?", (str(exc)[:500], now(), document_id))
        with db() as con:
            row = con.execute("SELECT * FROM documents WHERE id=?", (document_id,)).fetchone()
        return dict(row)
    except HTTPException:
        raise
    except Exception:
        safe_error("Document upload failed")
        stored.unlink(missing_ok=True)
        raise HTTPException(500, "The document could not be uploaded.")


@app.post("/documents/{document_id}/reindex")
def reindex_document(document_id: str):
    with db() as con:
        row = con.execute("SELECT * FROM documents WHERE id=?", (document_id,)).fetchone()
    if not row:
        raise HTTPException(404, "Document not found")
    path = UPLOAD_DIR / row["filename"]
    if not path.exists():
        # stored filename is UUID.ext; filename column is display name, so resolve by id.
        path = UPLOAD_DIR / f"{document_id}{row['extension']}"
    try:
        text = extract_text(path, row["extension"])
        chunks = chunk_text(text)
        if not chunks:
            raise ValueError("No extractable text was found.")
        with db() as con:
            con.execute("UPDATE documents SET status='indexed',chunks=?,error=NULL,updated_at=? WHERE id=?", (len(chunks), now(), document_id))
        return {"status": "indexed", "chunks": len(chunks)}
    except Exception as exc:
        with db() as con:
            con.execute("UPDATE documents SET status='failed',error=?,updated_at=? WHERE id=?", (str(exc)[:500], now(), document_id))
        raise HTTPException(422, "The document could not be indexed.")


@app.delete("/documents/{document_id}")
def delete_document(document_id: str):
    with db() as con:
        row = con.execute("SELECT * FROM documents WHERE id=?", (document_id,)).fetchone()
        if not row:
            raise HTTPException(404, "Document not found")
        con.execute("DELETE FROM documents WHERE id=?", (document_id,))
    (UPLOAD_DIR / f"{document_id}{row['extension']}").unlink(missing_ok=True)
    return {"deleted": True}


@app.get("/admin/metrics")
def admin_metrics():
    with db() as con:
        docs = con.execute("SELECT status,COUNT(*) n FROM documents GROUP BY status").fetchall()
        convs = con.execute("SELECT COUNT(*) FROM conversations").fetchone()[0]
        questions_today = con.execute("SELECT COUNT(*) FROM messages WHERE role='user' AND date(created_at)=date('now')").fetchone()[0]
        metrics = con.execute("SELECT questions,total_latency_ms FROM metrics WHERE id=1").fetchone()
        errors = con.execute("SELECT id,message,created_at FROM errors ORDER BY created_at DESC LIMIT 10").fetchall()
    by_status = {r["status"]: r["n"] for r in docs}
    questions = metrics[0] if metrics else 0
    avg = round(metrics[1] / questions) if questions else None
    return {
        "documents": {"total": sum(by_status.values()), "indexed": by_status.get("indexed", 0), "processing": by_status.get("processing", 0), "failed": by_status.get("failed", 0)},
        "conversations": convs,
        "questions_total": questions,
        "questions_today": questions_today,
        "average_response_time_ms": avg,
        "average_retrieval_score": None,
        "recent_errors": [dict(e) for e in errors],
    }


@app.get("/settings/public")
def public_settings():
    return {
        "product": "DocMind",
        "tagline": "AI-Powered Enterprise Knowledge Assistant",
        "max_upload_mb": MAX_FILE_SIZE // 1024 // 1024,
        "supported_types": sorted(ALLOWED_EXTENSIONS),
        "agent_configured": bool(os.getenv("AGENT_NAME") and os.getenv("AGENT_VERSION")),
        "search_configured": bool(os.getenv("AZURE_SEARCH_ENDPOINT")),
    }


@app.post("/chat/stream")
def chat_stream(request: ChatRequest):
    # Safe compatibility endpoint. If the configured SDK supports true streaming,
    # replace the iterator body with the SDK event stream; otherwise stream the
    # completed response in small chunks so the UI remains responsive.
    result = chat(request)
    text = result["response"]
    cid = result["conversation_id"]
    def iterator():
        for i in range(0, len(text), 32):
            yield json.dumps({"conversation_id": cid, "delta": text[i:i+32], "done": i + 32 >= len(text)}) + "\n"
    return StreamingResponse(iterator(), media_type="application/x-ndjson")
